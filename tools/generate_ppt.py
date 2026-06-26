#!/usr/bin/env python3
"""生成项目汇报 PPT — 4 页"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色 ──
DARK = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x00, 0x9f, 0xdb)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf0, 0xf0, 0xf5)
TEXT = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x2e, 0xcc, 0x71)
ORANGE = RGBColor(0xf3, 0x9c, 0x12)
RED = RGBColor(0xe7, 0x4c, 0x3c)
PURPLE = RGBColor(0x9b, 0x59, 0xb6)
GRAY_SUB = RGBColor(0x88, 0x88, 0x88)
DARK_BG_LEFT = RGBColor(0x22, 0x22, 0x3e)
DARK_BG_RIGHT = RGBColor(0x22, 0x3e, 0x22)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(tf, items, font_size=14, color=TEXT):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for par in cell.text_frame.paragraphs:
                par.font.size = Pt(11)
                par.font.name = "Microsoft YaHei"
                if r == 0:
                    par.font.bold = True
                    par.font.color.rgb = WHITE
                else:
                    par.font.color.rgb = TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table


def add_rounded_box(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, icon, title, items, title_color=ACCENT):
    """统一的功能卡片"""
    box = add_rounded_box(slide, left, top, width, height, fill_color=LIGHT_GRAY, border_color=RGBColor(0xdd, 0xdd, 0xdd))
    add_text_box(slide, left + 0.15, top + 0.1, width - 0.3, 0.35, f"{icon}  {title}", font_size=13, bold=True, color=title_color)
    tf = add_text_box(slide, left + 0.15, top + 0.5, width - 0.3, height - 0.6, "", font_size=10, color=TEXT)
    add_bullet_list(tf, items, font_size=10, color=TEXT)
    return box


# ═══════════════════════════════════════════
# Slide 1: 项目概述 — 两层结构 + MCP
# ═══════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, DARK)

# Title
add_text_box(slide1, 0.8, 0.3, 11.5, 0.8, "llm_wiki — 个人知识库 + 项目经验积累系统", font_size=36, bold=True, color=WHITE)
add_text_box(slide1, 0.8, 1.0, 11.5, 0.5, "AI 辅助知识管理 · 让经验不会随会话结束而丢失", font_size=18, color=ACCENT)

# ── 左：llm_wiki 底座 ──
add_text_box(slide1, 0.8, 1.8, 5.5, 0.45, "▎llm_wiki 知识库底座", font_size=20, bold=True, color=ACCENT)

base_shape = add_rounded_box(slide1, 0.8, 2.3, 5.5, 2.6, fill_color=DARK_BG_LEFT, border_color=ACCENT)

tf = add_text_box(slide1, 1.1, 2.4, 5.0, 2.3, "", font_size=14, color=WHITE)
add_bullet_list(tf, [
    "本地知识库管理 — 离线、私有、全本地",
    "全文 (BM25) + 向量 (ANN) 混合搜索，RRF 融合排序",
    "LLM 两阶段分析：源文件 → 分析 → 生成 → wiki 页面",
    "MCP Server — 8 个工具连接 Claude Code：",
    "  状态 / 项目列表 / 文件浏览 / 文件读取",
    "  搜索 / 知识图谱 / 审阅项 / 触发源扫描",
    "Tauri v2 + React 19 + LanceDB + Rust",
], font_size=14, color=WHITE)

# ── 中间连接 ──
add_text_box(slide1, 6.6, 3.0, 2.0, 1.0, "我在底座上\n搭建了 ↓", font_size=15, color=GRAY_SUB, alignment=PP_ALIGN.CENTER)

# ── 右：经验系统 ──
add_text_box(slide1, 7.2, 1.8, 5.5, 0.45, "▎项目经验积累系统", font_size=20, bold=True, color=GREEN)

exp_shape = add_rounded_box(slide1, 7.2, 2.3, 5.5, 2.6, fill_color=DARK_BG_RIGHT, border_color=GREEN)

tf = add_text_box(slide1, 7.5, 2.4, 5.0, 2.3, "", font_size=14, color=WHITE)
add_bullet_list(tf, [
    "会话结束自动捕获 → 提取经验 → 结构化入库",
    "Agent 启动自动加载项目情境（6 步）",
    "6 种经验类型：bug / decision / howto / ...",
    "错误发生时主动搜索并提醒",
    "/exp 随手标记，轻量不打断工作流",
    "Agent 自动配置 — 按文档逐步引导部署",
], font_size=14, color=WHITE)

# Bottom: 经验闭环流程
add_text_box(slide1, 0.8, 5.3, 12, 0.4, "▎经验闭环", font_size=18, bold=True, color=ACCENT)

steps = [
    ("写代码\n(Claude Code)", ACCENT),
    ("会话结束\n自动触发", RGBColor(0x34, 0x95, 0xdb)),
    ("提取经验\nLLM 分析", GREEN),
    ("入库存储\n可搜索", ORANGE),
    ("下次遇到\n自动提醒", RED),
]
for i, (label, color) in enumerate(steps):
    x = 1.4 + i * 2.3
    shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.85), Inches(1.9), Inches(0.85))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    if i < len(steps) - 1:
        arrow = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.95), Inches(6.1), Inches(0.28), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x55, 0x55, 0x55)
        arrow.line.fill.background()

add_text_box(slide1, 0.8, 6.9, 12, 0.4, "全自动闭环 · 完全本地化 · 支持 9 种 LLM Provider", font_size=12, color=GRAY_SUB, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 2: 核心功能（6 个功能卡片）
# ═══════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

add_text_box(slide2, 0.8, 0.3, 11.5, 0.7, "核心功能", font_size=32, bold=True, color=DARK)
add_text_box(slide2, 0.8, 0.85, 11.5, 0.4, "一次编码会话中，系统自动完成的工作", font_size=14, color=GRAY_SUB)

# Row 1 — 4 cards across top
cards_row1 = [
    {
        "icon": "📋",
        "title": "启动情境加载",
        "color": GREEN,
        "items": [
            "读取项目身份 (purpose.md)",
            "读取近期动态 (log.md)",
            "筛选未解决 bugs + 待定 decisions",
            "Git 上下文交叉比对",
            "→ 输出情境摘要",
        ],
    },
    {
        "icon": "🔍",
        "title": "错误主动搜索",
        "color": ACCENT,
        "items": [
            "编译报错 / 配置不生效 / 依赖缺失",
            "Agent 主动搜索知识库",
            "命中 → 提醒 + 摘要",
            "未命中 → 正常排查，解决后自动记录",
        ],
    },
    {
        "icon": "🏷️",
        "title": "/exp 随手标记",
        "color": ORANGE,
        "items": [
            "输入 /exp 刚才那个坑",
            "只输出几行轻量标记",
            "不写文件、不占上下文",
            "入库在会话结束后自动完成",
            "从「重操作」变「轻标记」",
        ],
    },
    {
        "icon": "📦",
        "title": "自动提取入库",
        "color": PURPLE,
        "items": [
            "SessionEnd Hook 自动触发",
            "JSONL → 过滤噪音 → 纯文本",
            "LLM 分析 → 生成经验页面",
            "向量嵌入 → LanceDB → 可搜索",
            "下次会话自动命中",
        ],
    },
]

for i, card in enumerate(cards_row1):
    x = 0.5 + i * 3.2
    add_card(slide2, x, 1.5, 3.0, 2.6, card["icon"], card["title"], card["items"], title_color=card["color"])

# Row 2 — 2 cards
cards_row2 = [
    {
        "icon": "🛡️",
        "title": "三层合并保护",
        "color": RED,
        "items": [
            "同一问题多次记录时不丢信息不重复",
            "第一层：sources/tags/related 自动去重并集（确定性）",
            "第二层：正文 LLM 智能整合（安全校验：不缩水 30%+）",
            "第三层：type/title/created 字段锁定不覆盖",
        ],
    },
    {
        "icon": "🗂️",
        "title": "6 种经验类型",
        "color": DARK,
        "items": [
            "bug — 缺陷记录，带状态流转 (unresolved→resolved)",
            "decision — 技术决策，方案对比 + 选择理由",
            "howto — 可复用操作步骤",
            "agent-error — AI 自身错误及纠正",
            "pattern — 重复坑，需≥2次证据才生成",
            "template — 可复用检查清单",
        ],
    },
]

for i, card in enumerate(cards_row2):
    x = 0.5 + i * 6.4
    add_card(slide2, x, 4.4, 6.0, 2.6, card["icon"], card["title"], card["items"], title_color=card["color"])

# ═══════════════════════════════════════════
# Slide 3: 如何使用
# ═══════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

add_text_box(slide3, 0.8, 0.3, 11.5, 0.7, "如何使用", font_size=32, bold=True, color=DARK)
add_text_box(slide3, 0.8, 0.85, 11.5, 0.4, "从部署到日常使用，只需三步", font_size=14, color=GRAY_SUB)

# Step 1: 部署
add_text_box(slide3, 0.8, 1.6, 2.0, 0.5, "第一步", font_size=14, bold=True, color=ACCENT)
add_text_box(slide3, 0.8, 2.0, 2.0, 0.4, "安装部署", font_size=24, bold=True, color=DARK)

step1_box = add_rounded_box(slide3, 0.5, 2.6, 3.8, 3.8, fill_color=LIGHT_GRAY)

tf = add_text_box(slide3, 0.8, 2.7, 3.3, 3.5, "", font_size=13, color=TEXT)
add_bullet_list(tf, [
    "① 安装 llm_wiki 应用",
    "  npm install + npm run tauri dev",
    "",
    "② 创建 wiki 项目",
    "  选择 Experience (🧠) 模板",
    "",
    "③ Agent 自动配置",
    "  把 agent-setup-guide.md 发给 AI",
    "  逐步引导完成全部配置",
    "",
    "④ 验证",
    "  检查 MCP 连通性",
    "  输入 /exp 测试",
], font_size=13, color=TEXT)

# Arrow 1→2
add_text_box(slide3, 4.5, 3.5, 1.0, 1.0, "→", font_size=36, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Step 2: 日常使用
add_text_box(slide3, 5.5, 1.6, 3.0, 0.5, "第二步", font_size=14, bold=True, color=GREEN)
add_text_box(slide3, 5.5, 2.0, 3.0, 0.4, "日常写代码", font_size=24, bold=True, color=DARK)

step2_box = add_rounded_box(slide3, 5.2, 2.6, 3.8, 3.8, fill_color=LIGHT_GRAY)

tf = add_text_box(slide3, 5.5, 2.7, 3.3, 3.5, "", font_size=13, color=TEXT)
add_bullet_list(tf, [
    "启动 Claude Code",
    "→ Agent 自动加载项目情境",
    "→ 告知未解决的 bugs / decisions",
    "",
    "开始写代码",
    "→ 错误发生 → 自动搜索经验",
    "→ 命中 → 提醒你以前踩过的坑",
    "",
    "想记录经验时",
    "→ 输入 /exp 刚才那个 xxx",
    "→ 只出一行标记，继续工作",
], font_size=13, color=TEXT)

# Arrow 2→3
add_text_box(slide3, 9.2, 3.5, 1.0, 1.0, "→", font_size=36, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# Step 3: 查看积累
add_text_box(slide3, 10.2, 1.6, 3.0, 0.5, "第三步", font_size=14, bold=True, color=PURPLE)
add_text_box(slide3, 10.2, 2.0, 3.0, 0.4, "查看积累", font_size=24, bold=True, color=DARK)

step3_box = add_rounded_box(slide3, 9.8, 2.6, 3.2, 3.8, fill_color=LIGHT_GRAY)

tf = add_text_box(slide3, 10.1, 2.7, 2.7, 3.5, "", font_size=13, color=TEXT)
add_bullet_list(tf, [
    "打开 llm_wiki 应用",
    "浏览经验页面",
    "",
    "按类型筛选：",
    "  bugs / decisions / patterns",
    "",
    "按状态筛选：",
    "  unresolved bugs",
    "  proposed decisions",
    "",
    "越用越多，持续积累",
], font_size=13, color=TEXT)

# Bottom: loop summary
add_text_box(slide3, 0.8, 6.7, 12, 0.5, "每次会话结束 → 经验自动入库 → 下次启动能搜到 → 越用越聪明", font_size=15, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 4: 现状与展望
# ═══════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)

add_text_box(slide4, 0.8, 0.3, 11.5, 0.7, "现状与展望", font_size=32, bold=True, color=DARK)

# ── 当前进展 ──
add_text_box(slide4, 0.8, 1.2, 5.5, 0.45, "▎当前进展", font_size=20, bold=True, color=DARK)

tf = add_text_box(slide4, 0.8, 1.7, 5.5, 1.8, "", font_size=13, color=TEXT)
add_bullet_list(tf, [
    "Windows 上完成模拟测试",
    "  自编 Claude Code 项目对话，核心功能基本跑通",
    "",
    "Linux 系统完成移植",
    "  代码已推送 GitHub，依赖已安装",
    "  修复了移植过程中的小问题",
], font_size=13, color=TEXT)

# ── 未解决的问题 ──
add_text_box(slide4, 0.8, 3.3, 11.5, 0.45, "▎未解决的问题", font_size=20, bold=True, color=DARK)

unsolved_data = [
    ["问题", "现状", "影响"],
    ["Linux SessionEnd Hook 未触发", "关闭会话窗口后没能自动导入经验到知识库，需要手动导入", "核心功能受影响"],
    ["大文件 Token 瓶颈", "超长会话（>300K 字符）可能导致生成截断，独立页面丢失", "长会话经验可能丢失"],
    ["多项目经验隔离", "结构支持了 project/domain，但搜索时能否自动按当前项目过滤未验证", "多项目共用时体验下降"],
    ["端到端测试不足", "模拟测试跑通，真实项目（121）未完整验证", "真实场景可能有意外"],
    ["[EXP] 标记端到端", "标记→识别→优先处理的完整链路未实测", "用户主动标记的可靠性不确定"],
]
add_table(slide4, 0.5, 3.8, 12.3, 2.1, 6, 3, unsolved_data, col_widths=[3.0, 6.3, 3.0])

# ── 下一步 ──
add_text_box(slide4, 0.8, 6.2, 11.5, 0.45, "▎下一步", font_size=18, bold=True, color=DARK)

tf = add_text_box(slide4, 0.8, 6.6, 11.5, 0.8, "", font_size=13, color=TEXT)
add_bullet_list(tf, [
    "修复 Linux SessionEnd Hook 触发 → 用 121 项目做端到端验证 → 解决多项目隔离 → 打包 .deb / .AppImage 发布",
], font_size=13, color=TEXT)

# ── Save ──
output_path = "docs/superpowers/项目汇报-经验积累系统-v2.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
